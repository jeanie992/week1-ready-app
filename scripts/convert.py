# -*- coding: utf-8 -*-
"""
convert.py  —  00_inbox 의 문서(PDF/Word/PPT/HWP)를 텍스트로 추출한다.

역할: '추출만' 한다. 분석/요약은 하지 않는다 (그건 analyze_excel.py 와 Claude /week1 담당).
입력: 00_inbox/pdf, 00_inbox/word, 00_inbox/ppt, 00_inbox/hwp
출력: 01_parsed/<원본명>.txt  (UTF-8)
      01_parsed/_manifest.csv  (파일명, 형식, 추출글자수, 상태, 경고)

설계 원칙:
- 파일 하나가 실패해도(try/except) 나머지는 계속 진행한다.
- 추출 글자수가 비정상적으로 적으면 경고를 남겨 사람이 확인하게 한다(스캔 PDF, HWP 표 누락 등).
- Excel/CSV 와 urls.txt 는 건드리지 않는다 (각각 analyze_excel.py, Claude 가 처리).
"""
from __future__ import annotations
import csv
import sys
from pathlib import Path

# ── 경로 설정 (이 파일 기준 상위 폴더가 프로젝트 루트) ──────────────────────
ROOT = Path(__file__).resolve().parent.parent
INBOX = ROOT / "00_inbox"
OUT = ROOT / "01_parsed"
MANIFEST = OUT / "_manifest.csv"

# 추출 글자수가 이 값 미만이면 "내용 거의 없음"으로 보고 경고
MIN_CHARS_WARN = 50


def extract_pdf(path: Path) -> str:
    """PDF → 텍스트. pdfplumber 사용(표 포함, MIT 라이선스)."""
    import pdfplumber
    parts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            txt = page.extract_text() or ""
            parts.append(f"\n--- [page {i}] ---\n{txt}")
    return "\n".join(parts).strip()


def extract_docx(path: Path) -> str:
    """Word(.docx) → 텍스트. 단락 + 표 셀 내용 포함."""
    import docx
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for ti, table in enumerate(d.tables, 1):
        parts.append(f"\n[표 {ti}]")
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts).strip()


def extract_pptx(path: Path) -> str:
    """PowerPoint(.pptx) → 텍스트. 슬라이드별 도형 텍스트."""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for si, slide in enumerate(prs.slides, 1):
        parts.append(f"\n--- [slide {si}] ---")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(c.text.strip() for c in row.cells))
    return "\n".join(parts).strip()


def extract_hwp(path: Path) -> str:
    """한글(.hwp/.hwpx) → 텍스트. hwpkit 사용(순수 파이썬, MIT)."""
    import hwpkit
    return (hwpkit.extract_text_from_file(str(path)) or "").strip()


# 형식별 (하위폴더, 확장자들, 추출함수) 매핑
HANDLERS = [
    ("pdf", {".pdf"}, extract_pdf),
    ("word", {".docx"}, extract_docx),
    ("ppt", {".pptx"}, extract_pptx),
    ("hwp", {".hwp", ".hwpx"}, extract_hwp),
]


def warn_for(kind: str, n_chars: int) -> str:
    """추출 결과에 대한 경고 메시지(없으면 빈 문자열)."""
    if n_chars == 0:
        if kind == "pdf":
            return "추출 0자 — 스캔본(이미지) 의심. OCR 또는 수동 확인 필요"
        if kind == "hwp":
            return "추출 0자 — 한글에서 PDF로 직접 내보내 00_inbox/pdf 에 넣으세요"
        return "추출 0자 — 파일 확인 필요"
    if n_chars < MIN_CHARS_WARN:
        if kind == "hwp":
            return "내용 매우 적음 — 표 누락 가능. 중요하면 PDF로 내보내 pdf 폴더에"
        return "내용 매우 적음 — 스캔본/빈 문서 여부 확인"
    return ""


def _clean_table(t):
    """노이즈(차트 텍스트박스) 제거 : 3행+·3열+·숫자 4개+·밀도 조건을 통과한 표만"""
    rows = [[(c or "").replace("\n", " ").strip() for c in r] for r in t]
    rows = [r for r in rows if any(r)]
    if len(rows) < 3:
        return None
    ncol = max(len(r) for r in rows)
    digits = sum(1 for r in rows for c in r if any(ch.isdigit() for ch in c))
    nonempty = sum(1 for r in rows for c in r if c)
    if ncol < 3 or digits < 4 or nonempty < len(rows) * 1.5:
        return None
    return rows


def extract_pdf_tables(path):
    """PDF 에서 '진짜 격자 표'만 추출(차트 이미지는 못 잡음) → (페이지, 표번호, 행들) 리스트"""
    import pdfplumber
    blocks = []
    with pdfplumber.open(path) as pdf:
        for i, pg in enumerate(pdf.pages, 1):
            for k, t in enumerate(pg.extract_tables() or [], 1):
                rows = _clean_table(t)
                if rows:
                    blocks.append((i, k, rows))
    return blocks


def tables_to_md(blocks, src):
    out = [f"# 추출된 표 : {src}",
           "> pdfplumber extract_tables : 격자 구조 표만(차트 이미지 제외) : 수치는 원본 PDF 대조 권장", ""]
    for i, k, rows in blocks:
        w = max(len(r) for r in rows)
        rows = [r + [""] * (w - len(r)) for r in rows]
        out.append(f"## p{i} 표{k}")
        out.append("| " + " | ".join(rows[0]) + " |")
        out.append("|" + "---|" * w)
        for r in rows[1:]:
            out.append("| " + " | ".join(r) + " |")
        out.append("")
    return "\n".join(out)


def main() -> int:
    if not INBOX.exists():
        print(f"[오류] 입력 폴더가 없습니다: {INBOX}", file=sys.stderr)
        return 1
    OUT.mkdir(parents=True, exist_ok=True)

    rows = []  # manifest 행
    total = ok = 0

    for kind, exts, fn in HANDLERS:
        folder = INBOX / kind
        if not folder.exists():
            continue
        for path in sorted(folder.iterdir()):
            if not path.is_file() or path.suffix.lower() not in exts:
                continue
            total += 1
            name = path.name
            try:
                text = fn(path)
                n = len(text)
                out_path = OUT / (path.stem + ".txt")
                # 출처 표시를 머리에 달아 Claude 가 인용하기 쉽게
                header = f"# 원본: {kind}/{name}\n\n"
                out_path.write_text(header + text, encoding="utf-8")
                ntab = 0
                if kind == "pdf":
                    try:
                        blocks = extract_pdf_tables(path)
                        ntab = len(blocks)
                        if blocks:
                            (OUT / (path.stem + "_tables.md")).write_text(tables_to_md(blocks, name), encoding="utf-8")
                    except Exception as te:
                        print(f"     (표 추출 건너뜀: {te})", file=sys.stderr)
                warn = warn_for(kind, n)
                status = "OK" if n >= MIN_CHARS_WARN else ("EMPTY" if n == 0 else "OK")
                if status == "OK":
                    ok += 1
                rows.append([name, kind, n, status, warn])
                tag = "OK " if not warn else "!! "
                tinfo = f", 표 {ntab}개" if ntab else ""
                print(f"  {tag}[{kind}] {name}  ({n:,}자{tinfo})  {warn}")
            except Exception as e:  # 한 파일 실패가 전체를 막지 않게
                rows.append([name, kind, 0, "ERROR", f"{type(e).__name__}: {e}"])
                print(f"  XX [{kind}] {name}  추출 실패: {e}", file=sys.stderr)

    # manifest 기록
    with MANIFEST.open("w", encoding="utf-8-sig", newline="") as f:
        w = csv.writer(f)
        w.writerow(["파일명", "형식", "추출글자수", "상태", "경고"])
        w.writerows(rows)

    print(f"\n완료: 총 {total}개 중 {ok}개 정상 추출. → {OUT}")
    print(f"상태표: {MANIFEST}")
    if total == 0:
        print("\n[안내] 00_inbox/{pdf,word,ppt,hwp} 에 자료를 먼저 넣어주세요.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
